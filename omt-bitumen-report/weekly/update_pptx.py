"""
update_pptx.py - Update PPTX tables with fresh data from collect_data.py.

Tables updated:
  Slide 2  "Таблица 6"  Russian NPZ prices
  Slide 3  "Таблица 11" ЖД deliveries (all countries)
  Slide 4  "Таблица 15" Currency rates
  Slide 4  "Таблица 19" Weather
  Slide 6  "Таблица 11" ЖД deliveries (Uzbekistan, 14 rows)
  Slide 7  "Таблица 4"  Uzbekistan price (сум/т)
  Slide 7  "Таблица 6"  Uzbekistan price ($/т)
  Slide 7  "Таблица 8"  Uzbekistan exchange weekly
  Slide 8  "Таблица 6"  ЖД deliveries (Uzbekistan, 5 rows)
  Slide 9  "Таблица 4"  ЖД deliveries (Kyrgyzstan)
  Slide 9  "Таблица 9"  ЖД deliveries (Other / Tajikistan)

Text boxes updated (analytical texts inserted separately via write_analytics.py):
  These shapes are left unchanged here.
"""
import sys
import math
import datetime
from pathlib import Path

from pptx import Presentation

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

# -------------------------------------------------------------------
# Formatting helpers
# -------------------------------------------------------------------

def _num(v, decimals=0):
    """Format number with space as thousands separator."""
    if v is None or (isinstance(v, float) and math.isnan(v)):
        return "-"
    if decimals == 0:
        return f"{round(v):,}".replace(",", "\u00a0")  # non-breaking space
    return f"{v:.{decimals}f}".replace(".", ",")


def _pct(v, decimals=0, zero_as_dash=True):
    """Format percent change value (already in 0..1 range)."""
    if v is None:
        return "-"
    pval = v * 100
    if zero_as_dash and abs(pval) < 0.05:
        return "-"
    sign = "+" if pval > 0 else ""
    if decimals == 0:
        return f"{sign}{round(pval)!s}%"
    return f"{sign}{pval:.{decimals}f}%".replace(".", ",")


def _pct_zhd(row):
    """Format ЖД % change: handles first-time delivery (prev=0, cur>0) as +100%."""
    cur  = row.get("vol_cur_month", 0) or 0
    prev = row.get("vol_prev_month", 0) or 0
    chg  = row.get("chg_month_pct")
    if prev == 0 and cur > 0:
        return "+100%"
    if chg is None:
        return "-"
    return _pct(chg)


def _temp(v):
    """Format temperature: +10,9 or -5,1"""
    if v is None:
        return "-"
    sign = "+" if v >= 0 else ""
    return f"{sign}{v:.1f}".replace(".", ",")


def _price(v):
    """Format ruble price: 26 800"""
    if v is None:
        return "-"
    return f"{round(v):,}".replace(",", "\u00a0")


def _price_usd(v):
    """Format USD price: integer"""
    if v is None:
        return "-"
    return str(round(v))


def _title_region(s):
    """Convert 'ОРЕНБУРГСКАЯ ОБЛАСТЬ' → 'Оренбургская область' (sentence case)."""
    if not s:
        return s
    # Proper nouns to keep capitalised mid-string
    PROPER = {"БАШКОРТОСТАН", "ТАТАРСТАН", "ДАГЕСТАН", "УДМУРТИЯ", "МОРДОВИЯ",
              "ЧУВАШИЯ", "МАРИЙЭЛ", "ХАКАСИЯ", "БУРЯТИЯ", "КОМИ", "ИНГУШЕТИЯ",
              "ЯКУТИЯ", "САНКТ-ПЕТЕРБУРГ", "МОСКВА", "МОСКОВСКАЯ"}
    words = s.split()
    result = []
    for i, w in enumerate(words):
        if i == 0 or w.upper() in PROPER:
            result.append(w.capitalize())
        else:
            result.append(w.lower())
    return " ".join(result)


def _fmt_date_range(dt_from, dt_to):
    """'с 23 по 27 марта 2026 г.'"""
    MONTHS = {
        1: "января", 2: "февраля", 3: "марта", 4: "апреля",
        5: "мая", 6: "июня", 7: "июля", 8: "августа",
        9: "сентября", 10: "октября", 11: "ноября", 12: "декабря",
    }
    if not isinstance(dt_from, datetime.datetime) or not isinstance(dt_to, datetime.datetime):
        return str(dt_from)
    m = MONTHS[dt_to.month]
    return f"с {dt_from.day} по {dt_to.day} {m} {dt_to.year} г."


# -------------------------------------------------------------------
# Cell setter (preserves font/fill, replaces text only)
# -------------------------------------------------------------------

def _set(cell, text):
    """Update cell text, preserving all formatting."""
    tf = cell.text_frame
    # Work on first paragraph, first run
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = str(text)
        for run in para.runs[1:]:
            run.text = ""
    else:
        para.add_run().text = str(text)
    # Clear extra paragraphs content
    for extra in tf.paragraphs[1:]:
        for run in extra.runs:
            run.text = ""


def _find_table(slide, name):
    for shape in slide.shapes:
        if shape.shape_type == 19 and shape.name == name:
            return shape.table
    return None


# -------------------------------------------------------------------
# Slide 2: Russian NPZ prices
# -------------------------------------------------------------------

def update_rf_prices(slide, rf_detail, report_date_str):
    tbl = _find_table(slide, "Таблица 6")
    if not tbl:
        print("  WARN: Slide 2 'Таблица 6' not found")
        return

    # Update date in header (row 0)
    if report_date_str:
        for ci in (3, 4, 6):
            cell = tbl.cell(0, ci)
            txt = cell.text_frame.text
            # Replace date pattern dd.mm.yyyy
            import re
            new_txt = re.sub(r"\d{2}\.\d{2}\.\d{4}", report_date_str, txt)
            if new_txt != txt:
                _set(cell, new_txt)

    # Fixed producer mapping
    from collect_data import RF_PRODUCERS
    total_row = None
    for k, v in rf_detail.items():
        if v.get("is_total"):
            total_row = v
            break

    for i, (region, company, npz_key) in enumerate(RF_PRODUCERS):
        row_idx = i + 1  # PPTX row (header is row 0)
        # Find matching entry in rf_detail
        found = None
        for (npz, comp), v in rf_detail.items():
            if v.get("is_total"):
                continue
            if npz_key.lower() in npz.lower() and company.lower()[:10] in comp.lower():
                found = v
                break

        if found:
            _set(tbl.cell(row_idx, 3), _price(found["price_s_nds"]))
            _set(tbl.cell(row_idx, 4), _price(found["price_b_nds"]))
            _set(tbl.cell(row_idx, 5), _pct(found["chg_rub"], decimals=1))
            _set(tbl.cell(row_idx, 6), _price_usd(found["price_usd"]))
            _set(tbl.cell(row_idx, 7), _pct(found["chg_usd"], decimals=1))
        else:
            # No data this week
            for ci in (3, 4, 5, 6, 7):
                _set(tbl.cell(row_idx, ci), "-*")

    # Total row (last)
    total_row_idx = tbl.rows.__len__() - 1
    if total_row:
        _set(tbl.cell(total_row_idx, 3), _price(total_row["price_s_nds"]))
        _set(tbl.cell(total_row_idx, 4), _price(total_row["price_b_nds"]))
        _set(tbl.cell(total_row_idx, 5), _pct(total_row["chg_rub"], decimals=1))
        _set(tbl.cell(total_row_idx, 6), _price_usd(total_row["price_usd"]))
        _set(tbl.cell(total_row_idx, 7), _pct(total_row["chg_usd"], decimals=1))

    print("  Slide 2 'Таблица 6' updated")


# -------------------------------------------------------------------
# Slide 3: ЖД all countries
# -------------------------------------------------------------------

def update_zhd_all(slide, zhd_all, month_ru):
    tbl = _find_table(slide, "Таблица 11")
    if not tbl:
        print("  WARN: Slide 3 'Таблица 11' not found")
        return

    # Update month in header row (row 0, col 3)
    if month_ru:
        import re
        header_cell = tbl.cell(0, 3)
        old_txt = header_cell.text_frame.text
        months_pattern = (
            "январ|феврал|март|апрел|май|июн|июл|август|сентябр|октябр|ноябр|декабр"
        )
        new_txt = re.sub(months_pattern, month_ru, old_txt, flags=re.IGNORECASE)
        if new_txt != old_txt:
            _set(header_cell, new_txt)

    n_data_rows = tbl.rows.__len__() - 2  # minus header and total
    # Separate total row
    total = next((r for r in zhd_all if r["region"] and "все страны" in r["region"].lower()), None)
    data_rows = [r for r in zhd_all if r["region"] and "все страны" not in r["region"].lower()]

    for i in range(n_data_rows):
        row_idx = i + 1
        if i < len(data_rows):
            r = data_rows[i]
            _set(tbl.cell(row_idx, 0), _title_region(r["region"]) if r["region"] else "-")
            _set(tbl.cell(row_idx, 1), r["company"])
            _set(tbl.cell(row_idx, 2), r["npz"])
            vol = r["vol_cur_month"]
            _set(tbl.cell(row_idx, 3), _num(vol) if vol else "-")
            _set(tbl.cell(row_idx, 4), _pct_zhd(r))
            _set(tbl.cell(row_idx, 5), _num(r["vol_cur_year"]) if r["vol_cur_year"] else "-")
        else:
            for ci in range(6):
                _set(tbl.cell(row_idx, ci), "-")

    # Total row
    total_row_idx = tbl.rows.__len__() - 1
    if total:
        _set(tbl.cell(total_row_idx, 3), _num(total["vol_cur_month"]))
        _set(tbl.cell(total_row_idx, 4), _pct(total["chg_month_pct"]))
        _set(tbl.cell(total_row_idx, 5), _num(total["vol_cur_year"]))

    print("  Slide 3 'Таблица 11' updated")


# -------------------------------------------------------------------
# Slide 4: Currency rates
# -------------------------------------------------------------------

def update_currency(slide, currency):
    tbl = _find_table(slide, "Таблица 15")
    if not tbl:
        print("  WARN: Slide 4 'Таблица 15' not found")
        return

    def _rate(v, decimals):
        if v is None:
            return "-"
        if decimals == 0:
            return f"{round(v):,}".replace(",", "\u00a0")
        return f"{v:.{decimals}f}".replace(".", ",")

    rows_data = [
        # (rate_value, decimals_for_rate, change_value)
        (currency["kzt"], 0, currency["kzt_chg"]),   # Казахстан тенге
        (currency["uzs"], 0, currency["uzs_chg"]),   # Узбекистан сум
        (currency["kgs"], 2, currency["kgs_chg"]),   # Кыргызстан сом
        (currency["tjs"], 1, currency["tjs_chg"]),   # Таджикистан сомони
        (currency["rub"], 2, currency["rub_chg"]),   # Россия рубль
    ]
    for i, (rate, dec, chg) in enumerate(rows_data):
        row_idx = i + 1
        _set(tbl.cell(row_idx, 1), _rate(rate, dec))
        _set(tbl.cell(row_idx, 2), _pct(chg, decimals=1))

    print("  Slide 4 'Таблица 15' updated")


# -------------------------------------------------------------------
# Slide 4: Weather
# -------------------------------------------------------------------

def update_weather(slide, weather):
    tbl = _find_table(slide, "Таблица 19")
    if not tbl:
        print("  WARN: Slide 4 'Таблица 19' not found")
        return

    for i, city in enumerate(weather):
        row_idx = i + 1
        if row_idx >= tbl.rows.__len__():
            break
        _set(tbl.cell(row_idx, 0), city["country"])
        _set(tbl.cell(row_idx, 1), city["city"])
        _set(tbl.cell(row_idx, 2), _temp(city["t_min"]))
        _set(tbl.cell(row_idx, 3), _temp(city["t_max"]))
        _set(tbl.cell(row_idx, 4), _temp(city["t_avg"]))

    print("  Slide 4 'Таблица 19' updated")


# -------------------------------------------------------------------
# Slides 6/8/9: ЖД per-country helper
# -------------------------------------------------------------------

def _update_zhd_country(slide, table_name, country_filter, zhd_per, tariffs, month_ru,
                         country_code=None, n_fixed_rows=None):
    """
    country_filter: callable(row) → bool, or string to match against row["country"]
    country_code: "Каз"/"Кирг"/"Узб" for tariff lookup
    n_fixed_rows: if given, table has exactly this many data rows (fixed layout)
    """
    tbl = _find_table(slide, table_name)
    if not tbl:
        print(f"  WARN: '{table_name}' not found")
        return

    if isinstance(country_filter, str):
        cf_str = country_filter.upper()
        country_filter = lambda r: r["country"] and cf_str in r["country"].upper()

    # Filter and sort
    rows = [r for r in zhd_per if country_filter(r)]
    # Separate total
    total = next((r for r in rows if "итого" in str(r["npz"]).lower()), None)
    data_rows = [r for r in rows if "итого" not in str(r["npz"]).lower()]

    has_tariff_col = tbl.rows[0].cells.__len__() >= 7

    # Update month header
    if month_ru:
        import re
        months_pattern = "январ|феврал|март|апрел|май|июн|июл|август|сентябр|октябр|ноябр|декабр"
        for ci in range(tbl.columns.__len__()):
            cell = tbl.cell(0, ci)
            old_txt = cell.text_frame.text
            new_txt = re.sub(months_pattern, month_ru, old_txt, flags=re.IGNORECASE)
            if new_txt != old_txt:
                _set(cell, new_txt)

    total_row_idx = tbl.rows.__len__() - 1
    n_data_rows = total_row_idx - 1  # rows 1 .. total_row_idx-1

    for i in range(n_data_rows):
        row_idx = i + 1
        if i < len(data_rows):
            r = data_rows[i]
            region  = _title_region(r["region"]) if r["region"] else "-"
            company = r["company"] or "-"
            npz     = r["npz"] or "-"
            vol     = r["vol_cur_month"]
            chg     = r["chg_month_pct"] if r["vol_prev_month"] else None
            yr_vol  = r["vol_cur_year"]

            # Tariff lookup
            tariff_str = "-"
            if has_tariff_col and country_code:
                tariff = tariffs.get((npz, country_code))
                if tariff:
                    tariff_str = _num(tariff, decimals=0)

            _set(tbl.cell(row_idx, 0), region)
            _set(tbl.cell(row_idx, 1), company)
            _set(tbl.cell(row_idx, 2), npz)
            _set(tbl.cell(row_idx, 3), _num(vol) if vol else "-")
            _set(tbl.cell(row_idx, 4), _pct_zhd(r))
            if has_tariff_col:
                _set(tbl.cell(row_idx, 5), tariff_str)
                _set(tbl.cell(row_idx, 6), _num(yr_vol) if yr_vol else "-")
            else:
                _set(tbl.cell(row_idx, 5), _num(yr_vol) if yr_vol else "-")
        else:
            n_cols = tbl.rows[0].cells.__len__()
            for ci in range(3, n_cols):
                _set(tbl.cell(row_idx, ci), "-")

    # Total row
    if total:
        vol   = total["vol_cur_month"]
        chg   = total["chg_month_pct"] if total["vol_prev_month"] else None
        yr    = total["vol_cur_year"]
        if has_tariff_col:
            # Average tariff for total row: keep existing or compute from filtered data
            avg_tariff_str = "-"
            valid = [(tariffs.get((r["npz"], country_code), 0) or 0) * (r["vol_cur_month"] or 0)
                     for r in data_rows if r["vol_cur_month"]]
            wsum  = sum(valid)
            wvol  = sum(r["vol_cur_month"] for r in data_rows if r["vol_cur_month"])
            if wvol > 0 and wsum > 0:
                avg_tariff_str = _num(wsum / wvol)
            _set(tbl.cell(total_row_idx, 3), _num(vol) if vol else "-")
            _set(tbl.cell(total_row_idx, 4), _pct(chg) if chg is not None else "-")
            _set(tbl.cell(total_row_idx, 5), avg_tariff_str)
            _set(tbl.cell(total_row_idx, 6), _num(yr) if yr else "-")
        else:
            _set(tbl.cell(total_row_idx, 3), _num(vol) if vol else "-")
            _set(tbl.cell(total_row_idx, 4), _pct(chg) if chg is not None else "-")
            _set(tbl.cell(total_row_idx, 5), _num(yr) if yr else "-")

    print(f"  '{table_name}' updated")


# -------------------------------------------------------------------
# Slide 7: Uzbekistan exchange prices
# -------------------------------------------------------------------

def update_uzb_prices(slide, uzb_exchange):
    # Таблица 4 (сум/т) and Таблица 6 ($/т)
    for tbl_name, price_key, chg_key in (
        ("Таблица 4", "price_sum", "chg_sum"),
        ("Таблица 6", "price_usd", "chg_usd"),
    ):
        tbl = _find_table(slide, tbl_name)
        if not tbl or not uzb_exchange:
            continue
        last = uzb_exchange[-1]
        if price_key == "price_sum":
            price_str = f"{round(last['price_sum']):,}".replace(",", "\u00a0")
        else:
            price_str = _price_usd(last["price_usd"])
        chg_str = _pct(last[chg_key], decimals=1)
        _set(tbl.cell(1, 1), price_str)
        _set(tbl.cell(1, 2), chg_str)
        print(f"  Slide 7 '{tbl_name}' updated")

    # Таблица 8: weekly exchange table (last 6 rows + delta)
    tbl8 = _find_table(slide, "Таблица 8")
    if tbl8 and uzb_exchange:
        n_data = tbl8.rows.__len__() - 2  # minus header and delta
        data_slice = uzb_exchange[-n_data:] if len(uzb_exchange) >= n_data else uzb_exchange
        for i, row in enumerate(data_slice):
            ri = i + 1
            date_label = _fmt_date_range(row["date_from"], row["date_to"])
            price_str  = f"{round(row['price_sum']):,}".replace(",", "\u00a0")
            vol_str    = str(round(row["vol"]))
            _set(tbl8.cell(ri, 0), date_label)
            _set(tbl8.cell(ri, 1), price_str)
            _set(tbl8.cell(ri, 2), vol_str)

        # Delta row (last row)
        delta_idx = tbl8.rows.__len__() - 1
        if len(uzb_exchange) >= 2:
            cur  = uzb_exchange[-1]
            prev = uzb_exchange[-2]
            dp   = round(cur["price_sum"]) - round(prev["price_sum"])
            dv   = round(cur["vol"]) - round(prev["vol"])
            _set(tbl8.cell(delta_idx, 1),
                 ("+" if dp >= 0 else "") + f"{dp:,}".replace(",", "\u00a0"))
            _set(tbl8.cell(delta_idx, 2),
                 ("+" if dv >= 0 else "") + str(dv))
        print("  Slide 7 'Таблица 8' updated")


# -------------------------------------------------------------------
# Main update entry point
# -------------------------------------------------------------------

def update_all(pptx_in: str, pptx_out: str, data: dict):
    print(f"\nОбновление таблиц в {Path(pptx_in).name}...")
    prs = Presentation(pptx_in)
    slides = prs.slides

    rf_detail     = data["rf_detail"]
    report_date_s = data["report_date_str"]
    month_ru      = data["month_ru"]
    weather       = data["weather"]
    currency      = data["currency"]
    zhd_all       = data["zhd_all"]
    zhd_per       = data["zhd_per"]
    tariffs       = data["tariffs"]
    uzb_exchange  = data["uzb_exchange"]

    # Slide indices (0-based)
    update_rf_prices(slides[1], rf_detail, report_date_s)   # Slide 2
    update_zhd_all(slides[2], zhd_all, month_ru)             # Slide 3
    update_currency(slides[3], currency)                      # Slide 4
    update_weather(slides[3], weather)                        # Slide 4
    _update_zhd_country(slides[5], "Таблица 11",             # Slide 6
                        "УЗБЕКИСТАН", zhd_per, tariffs, month_ru, "Узб")
    update_uzb_prices(slides[6], uzb_exchange)                # Slide 7
    _update_zhd_country(slides[7], "Таблица 6",              # Slide 8
                        "УЗБЕКИСТАН", zhd_per, tariffs, month_ru, "Узб")
    _update_zhd_country(slides[8], "Таблица 4",              # Slide 9
                        "КЫРГЫЗСТАН", zhd_per, tariffs, month_ru, "Кирг")
    _update_zhd_country(slides[8], "Таблица 9",              # Slide 9
                        lambda r: r["country"] and r["country"].upper() not in
                                  ("КАЗАХСТАН", "УЗБЕКИСТАН", "КЫРГЫЗСТАН"),
                        zhd_per, tariffs, month_ru, "Кирг")

    prs.save(pptx_out)
    print(f"\nСохранено: {pptx_out}")


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 3:
        print("Usage: python update_pptx.py <input.pptx> <output.pptx>")
        sys.exit(1)
    from collect_data import collect_all
    data = collect_all()
    update_all(sys.argv[1], sys.argv[2], data)
