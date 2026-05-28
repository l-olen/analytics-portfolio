"""
run_report.py - Weekly report automation entry point.

Usage:
  python run_report.py
      Collect data, search news + market prices, generate analytics, save to .txt files.

  python run_report.py --no-news
      Skip news and price search (faster, analytics based on data only).

  python run_report.py --write-pptx "Рынок СрАзии_13_04_2026.pptx"
      Also insert analytics texts into PPTX TextBoxes after generation.

Steps:
  1. Reads all Excel data (collect_data.py)
  2. Searches news via Tavily (search_news.py) + saves news_DD_MM_YYYY.txt
  3. Searches trader market prices via Tavily + saves market_prices_DD_MM_YYYY.txt
  4. Generates analytical texts via Claude API (generate_analytics.py)
  5. Saves texts to analytics_DD_MM_YYYY.txt
  6. [optional] Inserts texts into PPTX TextBoxes (write_analytics.py)
"""
import sys
import argparse
from pathlib import Path

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

BASE = Path(r"C:\projects\my-project\ОМТ\Еженедельный отчёт")


def _find_latest_file(pattern: str) -> str:
    """Return content of the most recently modified file matching glob pattern, or empty string."""
    files = sorted(BASE.glob(pattern), key=lambda p: p.stat().st_mtime, reverse=True)
    if files:
        print(f"  Загружаю: {files[0].name}")
        return files[0].read_text(encoding="utf-8")
    print(f"  Файл не найден: {pattern}")
    return ""


def main():
    parser = argparse.ArgumentParser(description="Weekly bitumen report analytics generator")
    parser.add_argument("--no-news", action="store_true",
                        help="Skip news and market price search (faster)")
    parser.add_argument("--analytics-only", action="store_true",
                        help="Skip Tavily search, use latest saved news/prices files")
    parser.add_argument("--write-pptx", metavar="FILE",
                        help="Insert analytics texts into this PPTX after generation")
    args = parser.parse_args()

    # Step 1: Read data from Excel
    from collect_data import collect_all
    data = collect_all()

    rdate = data.get("report_date")
    rdate = rdate.date() if hasattr(rdate, "date") else rdate

    # Step 2: Search news + market prices (or load from saved files)
    from search_news import (
        search_all_news, format_news_for_prompt, save_news,
        search_market_prices, format_prices_for_prompt, save_prices,
    )
    news_text = ""
    prices_text = ""

    if args.analytics_only:
        print("\nРежим --analytics-only: загружаю сохранённые файлы...")
        news_text = _find_latest_file("news_*.txt")
        prices_text = _find_latest_file("market_prices_*.txt")
    elif not args.no_news:
        print("\nПоиск новостей (последние 7 дней)...")
        news = search_all_news(verbose=True, days=7, report_date=rdate)
        save_news(news, report_date=rdate, out_dir=BASE)
        news_text = format_news_for_prompt(news)

        print("\nПоиск рыночных цен трейдеров (последние 14 дней)...")
        prices = search_market_prices(verbose=True, days=14, report_date=rdate)
        save_prices(prices, report_date=rdate, out_dir=BASE)
        prices_text = format_prices_for_prompt(prices)
    else:
        news_text = "Поиск новостей пропущен."

    # Step 2.5: Annotate global news with verified CA market impact
    if news_text and news_text.strip():
        print("\nВерификация глобальных новостей для рынка СА...")
        try:
            from news_verify import annotate_global_news
            news_text = annotate_global_news(news_text)
        except Exception as e:
            print(f"  [news-verify] Пропуск: {e}")

    # Step 2.6: Load previous analytics for style diff
    from generate_analytics import generate_texts, save_analytics, _load_prev_analytics
    cur_date_str = rdate.strftime("%d_%m_%Y") if rdate else ""
    prev_analytics = _load_prev_analytics(cur_date_str)
    if prev_analytics:
        print(f"  Предыдущий отчёт загружен для контроля стиля ({len(prev_analytics)} симв.)")
    else:
        print("  Предыдущий отчёт не найден, стилистический diff пропущен")

    # Step 3: Generate analytics via Claude
    analytics = generate_texts(data, news_text, prices_text, prev_text=prev_analytics)

    # Step 4: Save to file
    out = save_analytics(analytics, report_date=data.get("report_date"))

    # Step 5: Fact-check analytics against source data
    print("\nФакт-чекинг аналитики...")
    try:
        from fact_check_analytics import fact_check, save_fact_check
        from generate_analytics import _build_data_summary
        data_summary = _build_data_summary(data)
        fc_report = fact_check(analytics, data_summary, use_search=not args.no_news)
        save_fact_check(fc_report, report_date=data.get("report_date"))
        if fc_report["errors_found"]:
            print("  ! ВНИМАНИЕ: найдены расхождения с исходными данными -- проверь fact_check_*.txt")
    except Exception as e:
        print(f"  [fact-check] Ошибка: {e}")

    print(f"\nГотово. Файл: {out.name}")

    # Step 6 (optional): insert analytics into PPTX
    if args.write_pptx:
        from write_analytics import parse_analytics, find_latest, SLIDE_MAP, _find_shape, _write_to_textbox
        from pptx import Presentation

        pptx_path = BASE / args.write_pptx
        if not pptx_path.exists():
            print(f"\nWARN: PPTX не найден: {pptx_path}")
        else:
            print(f"\nВставка аналитики в {pptx_path.name}...")
            sections = parse_analytics(out.read_text(encoding="utf-8"))
            prs = Presentation(str(pptx_path))
            inserted = 0
            for slide_idx, shapes in SLIDE_MAP.items():
                if slide_idx >= len(prs.slides):
                    continue
                slide = prs.slides[slide_idx]
                for shape_name, section_key in shapes:
                    text = sections.get(section_key)
                    if not text:
                        continue
                    shape = _find_shape(slide, shape_name)
                    if not shape:
                        print(f"  WARN: слайд {slide_idx + 1}, '{shape_name}' не найден")
                        continue
                    _write_to_textbox(shape, text)
                    print(f"  OK: слайд {slide_idx + 1}, {shape_name} <- {section_key}")
                    inserted += 1
            prs.save(str(pptx_path))
            print(f"Сохранено: {pptx_path.name} ({inserted} блоков)")


if __name__ == "__main__":
    main()
