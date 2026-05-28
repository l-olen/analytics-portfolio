"""
write_analytics.py - Insert analytics texts into PPTX TextBoxes.

Reads the latest analytics_*.txt file, parses the 6 sections,
and writes each into the corresponding named TextBox in the PPTX.

TextBox mapping (from CLAUDE.md):
  Slide 1  TextBox14  ОБЗОР_РЫНКА
  Slide 2  TextBox10  РФ_НПЗ       (NOT TextBox20 -- that is a table caption)
  Slide 5  TextBox92  КАЗАХСТАН
  Slide 7  TextBox92  УЗБЕКИСТАН
  Slide 9  TextBox92  КЫРГЫЗСТАН
  Slide 9  TextBox21  ТАДЖИКИСТАН

Usage:
  python write_analytics.py                          -- latest analytics + latest PPTX
  python write_analytics.py --pptx "Рынок СрАзии_13_04_2026.pptx"
  python write_analytics.py --analytics analytics_10_04_2026.txt
  python write_analytics.py --dry-run                -- show parsed sections, no file changes
"""

import argparse
import re
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

BASE = Path(__file__).parent

# Slide index (0-based) → list of (shape_name, section_key)
SLIDE_MAP = {
    0: [("TextBox 14", "ОБЗОР_РЫНКА")],
    1: [("TextBox 10", "РФ_НПЗ")],
    4: [("TextBox 92", "КАЗАХСТАН")],
    6: [("TextBox 92", "УЗБЕКИСТАН")],
    8: [("TextBox 92", "КЫРГЫЗСТАН"), ("TextBox 21", "ТАДЖИКИСТАН")],
}

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ============================================================
# Parse analytics .txt file
# ============================================================
SECTION_HEADERS = {
    "обзор": "ОБЗОР_РЫНКА",
    "россия": "РФ_НПЗ",
    "нпз": "РФ_НПЗ",
    "казахстан": "КАЗАХСТАН",
    "узбекистан": "УЗБЕКИСТАН",
    "кыргызстан": "КЫРГЫЗСТАН",
    "таджикистан": "ТАДЖИКИСТАН",
}


def parse_analytics(text: str) -> dict[str, str]:
    """
    Parse analytics .txt into {section_key: text} dict.

    File format:
      РЫНОК БИТУМОВ...
      Дата: ...
      ====...====

      Обзор рынка
      -----------
      [text]

      ---

      Россия, НПЗ
      -----------
      [text]

      ---
      ...

    Sections are separated by blank line + --- + blank line.
    Header underlines (----) are longer and immediately follow the header -- NOT section separators.
    """
    # Strip file-level header (everything up to and including the ====... line)
    text = re.sub(r"^.*?={5,}\s*\n", "", text, flags=re.DOTALL)

    # Split by inter-section separator: blank line, ---, blank line
    parts = re.split(r"\n\n---\n\n", text)

    result = {}
    for part in parts:
        part = part.strip()
        if not part:
            continue

        lines = part.splitlines()
        if not lines:
            continue

        # First line is section header
        header = lines[0].strip()

        # Map header to section key
        section_key = None
        for keyword, key in SECTION_HEADERS.items():
            if keyword in header.lower():
                section_key = key
                break
        if not section_key:
            continue

        # Body starts after header (line 0) and optional underline (line 1 if all dashes)
        body_start = 1
        if len(lines) > 1 and re.match(r"^-+$", lines[1].strip()):
            body_start = 2

        body = "\n".join(lines[body_start:]).strip()
        if body:
            result[section_key] = body

    return result


# ============================================================
# Write text into PPTX TextBox
# ============================================================
def _find_shape(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name and shape.has_text_frame:
            return shape
    return None


HEADER_MAX_LEN = 60  # paragraphs shorter than this are treated as headers and preserved


def _para_text(para, tag_r: str, tag_t: str) -> str:
    """Extract plain text from a paragraph element."""
    return "".join(
        (r.find(tag_t).text or "")
        for r in para.findall(tag_r)
        if r.find(tag_t) is not None
    )


def _split_header_body(existing_paras: list, tag_r: str, tag_t: str):
    """
    Split existing paragraphs into (header_paras, body_paras).
    Header paragraphs are short (< HEADER_MAX_LEN chars) and come first.
    Example: ['Казахстан', '', 'Long analytics text...'] → header=['Казахстан',''], body=['Long...']
    """
    n = 0
    for para in existing_paras:
        txt = _para_text(para, tag_r, tag_t).strip()
        if len(txt) < HEADER_MAX_LEN:
            n += 1
        else:
            break
    return existing_paras[:n], existing_paras[n:]


def _find_body_template(body_paras: list, all_paras: list, tag_r: str, tag_t: str):
    """
    Find best paragraph to use as formatting template for body text.
    Prefers last body paragraph with runs.
    Falls back to last paragraph overall.
    Returns (template_para, template_run).
    """
    candidates = body_paras if body_paras else all_paras
    for para in reversed(candidates):
        runs = para.findall(tag_r)
        if runs:
            return para, runs[0]
    return all_paras[-1], None


def _clear_bold(run_elem, tag_rpr: str) -> None:
    """Remove bold attribute from run properties if present."""
    rpr = run_elem.find(tag_rpr)
    if rpr is not None and rpr.get("b") == "1":
        del rpr.attrib["b"]
    if rpr is not None and rpr.get("b") is not None:
        del rpr.attrib["b"]


def _write_to_textbox(shape, text: str) -> None:
    """
    Replace TextBox content with multi-paragraph text.
    Uses the LAST paragraph as formatting template to avoid copying
    bold/large header paragraphs (e.g. country name headers).
    Paragraphs in text are separated by double newlines.
    """
    from lxml import etree

    tf = shape.text_frame
    txBody = tf._txBody
    tag_p  = f"{{{NS}}}p"
    tag_r  = f"{{{NS}}}r"
    tag_t  = f"{{{NS}}}t"
    tag_rpr = f"{{{NS}}}rPr"

    # Split into non-empty paragraphs
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    if not paragraphs:
        return

    existing = txBody.findall(tag_p)
    if not existing:
        return

    # Preserve header paragraphs (short: section title, empty spacer)
    # Replace only body paragraphs
    header_paras, body_paras = _split_header_body(existing, tag_r, tag_t)
    template_para, template_run = _find_body_template(body_paras, existing, tag_r, tag_t)

    # Remove only body paragraphs (keep headers intact)
    for p in body_paras:
        txBody.remove(p)

    # Insert new paragraphs
    for para_text in paragraphs:
        new_p = deepcopy(template_para)

        # Clear existing runs from copy
        for r in new_p.findall(tag_r):
            new_p.remove(r)

        if template_run is not None:
            new_r = deepcopy(template_run)
            # Remove bold if present
            _clear_bold(new_r, tag_rpr)
            t_elem = new_r.find(tag_t)
            if t_elem is not None:
                t_elem.text = para_text
            else:
                t_elem = etree.SubElement(new_r, tag_t)
                t_elem.text = para_text
            new_p.append(new_r)
        else:
            # No runs in template -- create minimal run
            new_r = etree.SubElement(new_p, tag_r)
            t_elem = etree.SubElement(new_r, tag_t)
            t_elem.text = para_text

        txBody.append(new_p)


# ============================================================
# Main
# ============================================================
def find_latest(pattern: str) -> Path | None:
    files = sorted(BASE.glob(pattern), key=lambda p: p.stat().st_mtime, reverse=True)
    return files[0] if files else None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--pptx", help="Input PPTX file (default: latest in directory)")
    parser.add_argument("--analytics", help="Analytics .txt file (default: latest)")
    parser.add_argument("--output", help="Output PPTX file (default: overwrites input)")
    parser.add_argument("--dry-run", action="store_true", help="Parse and show sections only, no file changes")
    args = parser.parse_args()

    # Find analytics file
    analytics_path = Path(args.analytics) if args.analytics else find_latest("analytics_*.txt")
    if not analytics_path or not analytics_path.exists():
        print("Файл аналитики не найден. Укажите --analytics или запустите run_report.py.")
        sys.exit(1)

    print(f"Аналитика: {analytics_path.name}")
    sections = parse_analytics(analytics_path.read_text(encoding="utf-8"))

    if not sections:
        print("Не удалось распарсить секции из файла аналитики.")
        sys.exit(1)

    print(f"Найдено секций: {len(sections)} — {', '.join(sections.keys())}")

    if args.dry_run:
        print("\n--- DRY RUN: содержимое секций ---")
        for key, text in sections.items():
            print(f"\n[{key}]")
            print(text[:200] + "..." if len(text) > 200 else text)
        return

    # Find PPTX file
    pptx_path = Path(args.pptx) if args.pptx else find_latest("*.pptx")
    if not pptx_path or not pptx_path.exists():
        print("PPTX файл не найден. Укажите --pptx.")
        sys.exit(1)

    print(f"PPTX: {pptx_path.name}")
    output_path = Path(args.output) if args.output else pptx_path

    prs = Presentation(str(pptx_path))

    inserted = 0
    for slide_idx, shapes in SLIDE_MAP.items():
        if slide_idx >= len(prs.slides):
            print(f"  WARN: слайд {slide_idx + 1} не существует (всего {len(prs.slides)})")
            continue

        slide = prs.slides[slide_idx]
        for shape_name, section_key in shapes:
            text = sections.get(section_key)
            if not text:
                print(f"  SKIP: {section_key} — нет текста")
                continue

            shape = _find_shape(slide, shape_name)
            if not shape:
                print(f"  WARN: слайд {slide_idx + 1}, '{shape_name}' не найден")
                available = [s.name for s in slide.shapes if s.has_text_frame]
                print(f"        Доступные TextBox: {available}")
                continue

            tf = shape.text_frame
            paras_before = len(tf.paragraphs)
            runs_in_first = len(tf.paragraphs[0].runs) if tf.paragraphs else 0
            _write_to_textbox(shape, text)
            paras_after = len(tf.paragraphs)
            print(f"  OK:   слайд {slide_idx + 1}, {shape_name} <- {section_key} "
                  f"({len(text)} симв., параграфов: {paras_before}->{paras_after}, "
                  f"runs в 1м пар. до: {runs_in_first})")
            inserted += 1

    if inserted == 0:
        print("\nНичего не вставлено. Проверьте имена TextBox с помощью --dry-run и pptx_shapes.txt.")
        sys.exit(1)

    prs.save(str(output_path))
    print(f"\nСохранено: {output_path.name} ({inserted} блоков вставлено)")


if __name__ == "__main__":
    main()
