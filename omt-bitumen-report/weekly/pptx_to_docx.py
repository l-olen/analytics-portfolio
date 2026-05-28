"""
pptx_to_docx.py - Export analytics text boxes from weekly PPTX to Word.

Usage:
  python pptx_to_docx.py                        # latest PPTX in folder
  python pptx_to_docx.py --prev                 # second-to-last PPTX
  python pptx_to_docx.py "Рынок СрАзии_06_04_2026.pptx"
"""
import sys
import re
from pathlib import Path
from pptx import Presentation
from docx import Document
from docx.shared import Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH

BASE = Path(r"C:\projects\my-project\ОМТ\Еженедельный отчёт")

# (slide_index 0-based, shape_name, section_header)
TARGET_SHAPES = [
    (0, "TextBox 14", "Обзор рынка"),
    (1, "TextBox 10", "Россия, НПЗ"),
    (4, "TextBox 92", "Казахстан"),
    (6, "TextBox 92", "Узбекистан"),
    (8, "TextBox 92", "Кыргызстан"),
    (8, "TextBox 21", "Таджикистан"),
]


def find_pptx(path_arg=None, prev=False):
    if path_arg:
        p = Path(path_arg)
        if not p.is_absolute():
            p = BASE / p
        return p
    files = sorted(BASE.glob("Рынок СрАзии_*.pptx"), key=lambda f: f.stat().st_mtime)
    if not files:
        raise FileNotFoundError("PPTX файл не найден в папке")
    if prev:
        if len(files) < 2:
            raise FileNotFoundError("Предпоследний PPTX не найден (файлов меньше двух)")
        return files[-2]
    return files[-1]


def extract_text(prs, slide_idx, shape_name):
    slide = prs.slides[slide_idx]
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            return shape.text_frame.text.strip()
    return None


def clean_text(text, section_header):
    """Remove section header prefix if present at start of text box."""
    if text.startswith(section_header):
        text = text[len(section_header):].lstrip("\n ")
    return text


def build_docx(sections, pptx_name, out_path):
    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin    = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin   = Cm(2.5)
        section.right_margin  = Cm(2.5)

    # Title
    title = doc.add_heading(pptx_name.replace(".pptx", ""), level=1)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph()

    for header, text in sections:
        if not text:
            continue
        h = doc.add_heading(header, level=2)
        h.runs[0].font.size = Pt(13)

        p = doc.add_paragraph(text)
        p.paragraph_format.space_after = Pt(12)
        for run in p.runs:
            run.font.size = Pt(11)

    doc.save(str(out_path))


def main():
    prev = "--prev" in sys.argv
    path_arg = next((a for a in sys.argv[1:] if a != "--prev"), None)
    pptx_path = find_pptx(path_arg, prev=prev)
    print(f"Читаю: {pptx_path.name}")

    prs = Presentation(str(pptx_path))

    sections = []
    for slide_idx, shape_name, header in TARGET_SHAPES:
        text = extract_text(prs, slide_idx, shape_name)
        if text:
            text = clean_text(text, header)
            sections.append((header, text))
            print(f"  + {header} ({len(text)} символов)")
        else:
            print(f"  ! {header} - не найдено (слайд {slide_idx+1}, {shape_name})")

    out_path = pptx_path.with_suffix(".docx")
    build_docx(sections, pptx_path.name, out_path)
    print(f"\nСохранено: {out_path.name}")


if __name__ == "__main__":
    main()
